from bs4 import BeautifulSoup
from playwright.sync_api import sync_playwright
from pypdf import PdfWriter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pdf2image import convert_from_path
import os
from PIL import Image
import io
import re

# HTML 파일 경로
html_file = 'figma_import_version.html'

# HTML 파일 로드 및 파싱
with open(html_file, 'r', encoding='utf-8') as f:
    html_content = f.read()
    soup = BeautifulSoup(html_content, 'html.parser')

# 원본 head 부분을 포맷팅 유지하며 추출 (정규표현식 사용)
head_match = re.search(r'<head[^>]*>.*?</head>', html_content, re.DOTALL | re.IGNORECASE)
if head_match:
    head_content = head_match.group(0)
else:
    # 정규표현식 실패 시 BeautifulSoup 사용 (fallback)
    head_tag = soup.find('head')
    head_content = str(head_tag) if head_tag else ''

# head에서 스타일 정보 추출
style_tag = soup.find('style')
css_rules = {}
if style_tag:
    # CSS 규칙 파싱 (간단한 버전)
    css_text = style_tag.get_text()
    # .slide { ... } 형식 파싱
    for rule in re.finditer(r'\.(\w+)\s*\{([^}]+)\}', css_text):
        class_name = rule.group(1)
        properties = rule.group(2)
        css_rules[class_name] = {}
        for prop in re.finditer(r'([^:;]+):\s*([^;]+);', properties):
            key = prop.group(1).strip()
            value = prop.group(2).strip()
            css_rules[class_name][key] = value

# Tailwind 색상 설정 추출 (wow.html의 tailwind.config에서)
tailwind_colors = {}
tailwind_config_script = soup.find('script', string=re.compile(r'tailwind\.config'))
if tailwind_config_script:
    config_text = tailwind_config_script.string
    # colors 객체에서 색상 추출
    colors_match = re.search(r'colors:\s*\{([^}]+)\}', config_text, re.DOTALL)
    if colors_match:
        colors_text = colors_match.group(1)
        # 'color-name': '#HEX' 형식 추출
        color_matches = re.finditer(r"['\"]([^'\"]+)['\"]:\s*['\"](#[0-9A-Fa-f]{6})['\"]", colors_text)
        for match in color_matches:
            color_name = match.group(1)
            color_hex = match.group(2)
            tailwind_colors[color_name] = color_hex

# 기본값 설정 (추출 실패 시)
if not tailwind_colors:
    tailwind_colors = {
        'soft-white': '#F9F8F6',
        'dark-gray': '#333333',
        'coral-red': '#FF6B6B',
        'coral-dark': '#E05252',
        'light-gray': '#EAEAEA',
    }

# 슬라이드 추출 (section 태그 기준)
slides = soup.find_all('section')

# 원본 body 구조 확인
body_tag = soup.find('body')
body_classes = body_tag.get('class', []) if body_tag else []
body_class_str = ' '.join(body_classes) if body_classes else ''

# 원본 presentation-container 구조 확인
presentation_container = soup.find('div', id='presentation-container')
presentation_container_classes = presentation_container.get('class', []) if presentation_container else []
presentation_container_class_str = ' '.join(presentation_container_classes) if presentation_container_classes else ''

# 원본 slides-wrapper 구조 확인
slides_wrapper = soup.find('div', id='slides-wrapper')
slides_wrapper_classes = slides_wrapper.get('class', []) if slides_wrapper else []
slides_wrapper_class_str = ' '.join(slides_wrapper_classes) if slides_wrapper_classes else 'w-full h-full relative'
slides_wrapper_style = slides_wrapper.get('style', '') if slides_wrapper else ''

# 임시 슬라이드 html 저장 경로 리스트
slide_html_files = []

for idx, slide in enumerate(slides):
    # 각 슬라이드에 원본 head를 그대로 사용하고, 원본 구조를 완전히 유지
    # 단, 개별 슬라이드일 때는 position: absolute가 문제가 되므로 relative로 변경하는 CSS 추가
    
    # section에 active 클래스 추가 (원본 CSS에서 .slide.active만 보이도록 설정되어 있음)
    slide_classes = slide.get('class', [])
    if 'active' not in slide_classes:
        slide_classes.append('active')
    slide['class'] = slide_classes
    
    # 개별 슬라이드를 위한 CSS 오버라이드 추가 (position: absolute -> relative로 변경)
    slide_specific_css = """
    <style>
        /* 개별 슬라이드일 때 position: absolute가 레이아웃을 깨뜨리므로 relative로 변경 */
        .slide {
            position: relative !important;
        }
        /* slides-wrapper도 relative로 확실히 설정 */
        #slides-wrapper {
            position: relative !important;
        }
    </style>
    """
    
    # head에 슬라이드별 CSS 추가
    head_with_slide_css = head_content.replace('</head>', slide_specific_css + '</head>')
    
    # wow.html의 원본 구조 완전히 유지: body > presentation-container > slides-wrapper > section
    slide_html = f"""<!DOCTYPE html>
<html lang="ko">
{head_with_slide_css}
<body class="{body_class_str}">
    <div id="presentation-container" class="{presentation_container_class_str}">
        <div id="slides-wrapper" class="{slides_wrapper_class_str}" style="{slides_wrapper_style}">
            {slide}
        </div>
    </div>
</body>
</html>"""

    temp_html_path = f"slide_{idx+1}.html"
    # 절대 경로로 변환
    abs_path = os.path.abspath(temp_html_path)
    with open(abs_path, 'w', encoding='utf-8') as tmp:
        tmp.write(slide_html)
    
    slide_html_files.append(abs_path)

print(f"총 {len(slide_html_files)}개의 슬라이드 HTML 파일 생성 완료")

# Playwright를 이용한 PDF 생성
pdf_paths = []  # 생성된 PDF 경로 저장

with sync_playwright() as p:
    browser = p.chromium.launch()
    # 원본 HTML에 맞게 자연스럽게 렌더링 (고정 viewport 제거)
    page = browser.new_page()

    for idx, html_path in enumerate(slide_html_files):
        file_url = f"file://{html_path}"
        
        print(f"슬라이드 {idx+1} 처리 중...")
        page.goto(file_url, wait_until='networkidle')
        
        # 모든 이미지와 외부 리소스가 완전히 로드될 때까지 대기
        page.wait_for_load_state('networkidle')
        
        # 배경 이미지가 로드될 때까지 추가 대기
        page.wait_for_timeout(2000)
        
        # 모든 이미지가 로드되었는지 확인
        images_loaded = page.evaluate("""
            () => {
                return Promise.all(
                    Array.from(document.images).map(img => {
                        if (img.complete) return Promise.resolve();
                        return new Promise((resolve, reject) => {
                            img.onload = resolve;
                            img.onerror = resolve; // 에러가 나도 계속 진행
                            setTimeout(resolve, 3000); // 최대 3초 대기
                        });
                    })
                );
            }
        """)
        
        # CSS 배경 이미지도 로드될 때까지 추가 대기
        page.wait_for_timeout(1000)
        
        # presentation-container의 실제 크기 측정
        container_size = page.evaluate("""
            () => {
                const container = document.getElementById('presentation-container');
                if (container) {
                    const rect = container.getBoundingClientRect();
                    return {
                        width: Math.ceil(rect.width),
                        height: Math.ceil(rect.height)
                    };
                }
                return { width: 1920, height: 1080 }; // 기본값
            }
        """)
        
        # PDF 저장 (원본 HTML의 실제 렌더링 크기에 맞춤)
        pdf_path = f"slide_{idx+1}.pdf"
        # 절대 경로로 변환
        abs_pdf_path = os.path.abspath(pdf_path)
        
        page.pdf(
            path=abs_pdf_path,
            width=f"{container_size['width']}px",
            height=f"{container_size['height']}px",
            print_background=True,
            margin={"top": "0", "right": "0", "bottom": "0", "left": "0"}
        )

        # PDF 파일이 제대로 생성되었는지 확인
        if os.path.exists(abs_pdf_path) and os.path.getsize(abs_pdf_path) > 0:
            pdf_paths.append(abs_pdf_path)
            print(f"✓ {pdf_path} 생성 완료")
        else:
            print(f"⚠ {pdf_path} 생성 실패 또는 파일이 비어있음")

    browser.close()

print(f"\n개별 PDF 생성 완료! 총 {len(pdf_paths)}개 파일")

# 모든 PDF를 하나로 병합
print("\nPDF 병합 중...")
merger = PdfWriter()

for pdf_path in pdf_paths:
    # PDF 파일이 존재하고 유효한지 확인
    if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
        try:
            merger.append(pdf_path)
            print(f"  - {os.path.basename(pdf_path)} 추가됨")
        except Exception as e:
            print(f"  ⚠ {os.path.basename(pdf_path)} 병합 실패: {e}")
    else:
        print(f"  ⚠ {os.path.basename(pdf_path)} 파일이 존재하지 않거나 비어있음")

# 최종 병합된 PDF 저장
output_pdf = "merged_slides.pdf"
merger.write(output_pdf)
merger.close()

print(f"\n✓ 최종 결과물: {output_pdf} 생성 완료!")
print(f"  총 {len(pdf_paths)}개 슬라이드가 하나의 PDF로 병합되었습니다.")

# HTML을 PPTX로 직접 변환 (편집 가능한 객체로)
print("\n" + "="*50)
print("HTML을 PPTX로 변환 중 (편집 가능한 객체로)...")
print("="*50)

def parse_color(color_str):
    """색상 문자열을 RGBColor로 변환"""
    if not color_str:
        return None
    
    # #005A9C 형식
    if color_str.startswith('#'):
        try:
            r = int(color_str[1:3], 16)
            g = int(color_str[3:5], 16)
            b = int(color_str[5:7], 16)
            return RGBColor(r, g, b)
        except:
            pass
    
    # rgb() 형식
    rgb_match = re.search(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', color_str)
    if rgb_match:
        return RGBColor(int(rgb_match.group(1)), int(rgb_match.group(2)), int(rgb_match.group(3)))
    
    # 색상 이름 매핑
    color_map = {
        'slide-blue': RGBColor(0, 90, 156),  # #005A9C
        'slide-text': RGBColor(33, 37, 41),   # #212529
        'slide-gray': RGBColor(108, 117, 125), # #6C757D
    }
    return color_map.get(color_str.lower())

def get_font_size(class_str):
    """클래스에서 폰트 크기 추출"""
    if 'text-4xl' in class_str or 'text-5xl' in class_str:
        return Pt(44)
    elif 'text-3xl' in class_str:
        return Pt(36)
    elif 'text-2xl' in class_str:
        return Pt(28)
    elif 'text-xl' in class_str:
        return Pt(24)
    elif 'text-lg' in class_str:
        return Pt(20)
    elif 'text-sm' in class_str:
        return Pt(14)
    return Pt(18)  # 기본값

def is_bold(class_str):
    """클래스에서 bold 여부 확인"""
    return 'font-bold' in class_str or 'font-semibold' in class_str

def clean_text(text):
    """텍스트에서 마크다운 형식 제거"""
    # **텍스트** -> 텍스트 (볼드)
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    return text.strip()

def html_to_pptx_slide(slide_soup, prs, slide_idx, css_rules, tailwind_colors):
    """HTML 슬라이드를 PPTX 슬라이드로 변환 (head의 스타일 정보 반영)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    
    # 슬라이드 배경색 설정 (CSS에서 .slide의 background-color 또는 인라인 스타일)
    slide_style = slide_soup.get('style', '')
    bg_color = RGBColor(248, 249, 250)  # 기본값 #F8F9FA
    
    # 인라인 스타일에서 배경색 추출
    if 'background' in slide_style:
        bg_match = re.search(r'background[^:]*:\s*([^;]+)', slide_style)
        if bg_match:
            bg_value = bg_match.group(1)
            # linear-gradient는 첫 번째 색상 사용
            color_match = re.search(r'#([0-9A-Fa-f]{6})', bg_value)
            if color_match:
                bg_color = parse_color('#' + color_match.group(1))
    
    # CSS 규칙에서 .slide의 background-color 확인
    if 'slide' in css_rules and 'background-color' in css_rules['slide']:
        bg_color = parse_color(css_rules['slide']['background-color']) or bg_color
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # 패딩 계산 (CSS에서 .slide-header, .slide-body의 padding 참고)
    header_padding_top = Inches(2.5 / 16)  # 2.5rem ≈ 2.5/16 인치
    header_padding_horizontal = Inches(3 / 16)  # 3rem
    body_padding_horizontal = Inches(3 / 16)
    body_padding_top = Inches(1.5 / 16)
    body_padding_bottom = Inches(3 / 16)
    
    if 'slide-header' in css_rules:
        padding = css_rules['slide-header'].get('padding', '')
        if padding:
            # padding: 2.5rem 3rem 0 3rem 형식 파싱
            padding_parts = re.findall(r'([\d.]+)rem', padding)
            if len(padding_parts) >= 1:
                header_padding_top = Inches(float(padding_parts[0]) / 16)
            if len(padding_parts) >= 2:
                header_padding_horizontal = Inches(float(padding_parts[1]) / 16)
    
    y_position = header_padding_top
    left_margin = header_padding_horizontal
    right_margin = prs.slide_width - header_padding_horizontal
    content_width = right_margin - left_margin
    
    # slide-header 처리
    header = slide_soup.find('div', class_='slide-header')
    if header:
        h2 = header.find('h2')
        if h2:
            text = clean_text(h2.get_text())
            # CSS에서 h2 스타일 적용
            h2_font_size = Pt(32)  # 2rem = 32pt
            h2_color = RGBColor(33, 37, 41)  # #212529
            if 'h2' in css_rules:
                if 'font-size' in css_rules['h2']:
                    size_match = re.search(r'([\d.]+)rem', css_rules['h2']['font-size'])
                    if size_match:
                        h2_font_size = Pt(float(size_match.group(1)) * 16)
                if 'color' in css_rules['h2']:
                    h2_color = parse_color(css_rules['h2']['color']) or h2_color
            
            textbox = slide.shapes.add_textbox(left_margin, y_position, content_width, Inches(1.2))
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = text
            p.font.size = h2_font_size
            p.font.bold = True
            p.font.color.rgb = h2_color
            
            # h2의 border-bottom 추가 (CSS에서 4px solid #005A9C)
            border_y = y_position + Inches(1.0)
            line = slide.shapes.add_connector(1, left_margin, border_y, left_margin + Inches(2), border_y)
            line.line.color.rgb = parse_color('#005A9C')
            line.line.width = Pt(4)
            
            y_position += Inches(1.5)
    
    # slide-body 처리 (body 패딩 적용)
    body = slide_soup.find('div', class_='slide-body')
    if body:
        body_x = left_margin + body_padding_horizontal
        body_y = y_position + body_padding_top
        body_width = content_width - (body_padding_horizontal * 2)
        
        # 그리드 레이아웃 확인
        grid = body.find('div', class_=re.compile(r'grid'))
        if grid:
            cols = 2 if 'grid-cols-2' in str(grid.get('class', [])) else 3 if 'grid-cols-3' in str(grid.get('class', [])) else 1
            col_width = body_width / cols
            col_gap = Inches(0.5)
            
            items = grid.find_all(['div'], recursive=False)
            current_col = 0
            current_y = body_y
            
            for item in items:
                col_x = body_x + (col_width + col_gap) * current_col
                process_element(item, slide, col_x, current_y, col_width - col_gap, prs, css_rules, tailwind_colors)
                current_col += 1
                if current_col >= cols:
                    current_col = 0
                    current_y += Inches(3)  # 다음 행
        else:
            # 일반 레이아웃
            process_element(body, slide, body_x, body_y, body_width, prs, css_rules, tailwind_colors)
    
    # slide-footer 처리
    footer = slide_soup.find('div', class_='slide-footer')
    if footer:
        # CSS에서 .slide-footer의 padding, font-size, color 확인
        footer_padding = Inches(0.75 / 16)  # 0.75rem
        footer_font_size = Pt(14)  # 0.875rem
        footer_color = RGBColor(108, 117, 125)  # #6C757D
        
        if 'slide-footer' in css_rules:
            if 'font-size' in css_rules['slide-footer']:
                size_match = re.search(r'([\d.]+)rem', css_rules['slide-footer']['font-size'])
                if size_match:
                    footer_font_size = Pt(float(size_match.group(1)) * 16)
            if 'color' in css_rules['slide-footer']:
                footer_color = parse_color(css_rules['slide-footer']['color']) or footer_color
        
        footer_y = prs.slide_height - Inches(0.75)
        footer_text = clean_text(footer.get_text())
        textbox = slide.shapes.add_textbox(left_margin, footer_y, content_width, Inches(0.4))
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = footer_text
        p.font.size = footer_font_size
        p.font.color.rgb = footer_color
        p.alignment = PP_ALIGN.JUSTIFY
        
        # border-top 추가
        border_y = footer_y
        line = slide.shapes.add_connector(1, left_margin, border_y, right_margin, border_y)
        line.line.color.rgb = parse_color('#DEE2E6')  # slide-border
        line.line.width = Pt(1)

def process_element(element, slide, x, y, width, prs, css_rules, tailwind_colors, depth=0):
    """HTML 요소를 재귀적으로 처리하여 PPTX 객체로 변환"""
    if depth > 10:  # 깊이 제한
        return y
    
    current_y = y
    
    # h1, h2, h3 처리
    for tag in ['h1', 'h2', 'h3']:
        headings = element.find_all(tag, recursive=False)
        for h in headings:
            text = clean_text(h.get_text())
            if text:
                class_str = ' '.join(h.get('class', []))
                
                # CSS 규칙에서 태그 스타일 확인
                font_size = get_font_size(class_str)
                font_color = parse_color('slide-blue') or RGBColor(0, 90, 156)
                is_bold_text = is_bold(class_str)
                
                if tag in css_rules:
                    if 'font-size' in css_rules[tag]:
                        size_match = re.search(r'([\d.]+)rem', css_rules[tag]['font-size'])
                        if size_match:
                            font_size = Pt(float(size_match.group(1)) * 16)
                    if 'color' in css_rules[tag]:
                        font_color = parse_color(css_rules[tag]['color']) or font_color
                    if 'font-weight' in css_rules[tag]:
                        is_bold_text = '700' in css_rules[tag]['font-weight'] or '600' in css_rules[tag]['font-weight']
                
                # Tailwind 클래스에서 색상 확인
                for color_name, color_hex in tailwind_colors.items():
                    if f'text-{color_name}' in class_str or f'bg-{color_name}' in class_str:
                        font_color = parse_color(color_hex)
                        break
                
                textbox = slide.shapes.add_textbox(x, current_y, width, Inches(0.8))
                text_frame = textbox.text_frame
                text_frame.word_wrap = True
                p = text_frame.paragraphs[0]
                p.text = text
                p.font.size = font_size
                p.font.bold = is_bold_text
                p.font.color.rgb = font_color
                current_y += Inches(0.9)
    
    # ul, ol 리스트 처리
    lists = element.find_all(['ul', 'ol'], recursive=False)
    for ul in lists:
        # CSS에서 ul 스타일 확인
        list_font_size = Pt(18)  # 1.125rem
        list_color = RGBColor(52, 58, 79)  # #343A4F
        
        if 'ul' in css_rules:
            if 'font-size' in css_rules['ul']:
                size_match = re.search(r'([\d.]+)rem', css_rules['ul']['font-size'])
                if size_match:
                    list_font_size = Pt(float(size_match.group(1)) * 16)
            if 'color' in css_rules['ul']:
                list_color = parse_color(css_rules['ul']['color']) or list_color
        
        items = ul.find_all('li', recursive=False)
        for li in items:
            text = clean_text(li.get_text())
            if text:
                # CSS에서 li::before의 content 확인 (■)
                bullet = "■"
                if 'ul > li' in css_rules or 'li::before' in str(css_rules):
                    # CSS에서 content 확인
                    pass  # 기본값 사용
                
                textbox = slide.shapes.add_textbox(x + Inches(0.3), current_y, width - Inches(0.3), Inches(0.6))
                text_frame = textbox.text_frame
                text_frame.word_wrap = True
                p = text_frame.paragraphs[0]
                p.text = f"{bullet} {text}"
                p.font.size = list_font_size
                p.font.color.rgb = list_color
                current_y += Inches(0.7)
    
    # p 태그 처리
    paragraphs = element.find_all('p', recursive=False)
    for p_tag in paragraphs:
        text = clean_text(p_tag.get_text())
        if text:
            class_str = ' '.join(p_tag.get('class', []))
            
            # CSS에서 p 스타일 확인
            p_font_size = Pt(18)  # 1.125rem
            p_color = RGBColor(52, 58, 79)  # #343A4F
            
            if 'p' in css_rules:
                if 'font-size' in css_rules['p']:
                    size_match = re.search(r'([\d.]+)rem', css_rules['p']['font-size'])
                    if size_match:
                        p_font_size = Pt(float(size_match.group(1)) * 16)
                if 'color' in css_rules['p']:
                    p_color = parse_color(css_rules['p']['color']) or p_color
            
            # Tailwind 클래스 우선 적용
            p_font_size = get_font_size(class_str) or p_font_size
            for color_name, color_hex in tailwind_colors.items():
                if f'text-{color_name}' in class_str:
                    p_color = parse_color(color_hex)
                    break
            
            textbox = slide.shapes.add_textbox(x, current_y, width, Inches(0.5))
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = text
            p.font.size = p_font_size
            p.font.color.rgb = p_color
            current_y += Inches(0.6)
    
    # div 재귀 처리 (중첩된 구조)
    divs = element.find_all('div', recursive=False)
    for div in divs:
        # 특정 클래스는 건너뛰기
        div_class = ' '.join(div.get('class', []))
        if 'figure-placeholder' in div_class:
            continue
        current_y = process_element(div, slide, x, current_y, width, prs, css_rules, tailwind_colors, depth + 1)
    
    return current_y

# PPTX 프레젠테이션 생성 (16:9 비율)
prs = Presentation()
prs.slide_width = Inches(10)  # 16:9 비율의 너비
prs.slide_height = Inches(5.625)  # 16:9 비율의 높이

print(f"\n총 {len(slides)}개 슬라이드를 PPTX로 변환 중...")

for idx, slide_html in enumerate(slides):
    print(f"슬라이드 {idx+1}/{len(slides)} 처리 중...")
    try:
        html_to_pptx_slide(slide_html, prs, idx, css_rules, tailwind_colors)
        print(f"  ✓ 슬라이드 {idx+1} 변환 완료")
    except Exception as e:
        print(f"  ⚠ 슬라이드 {idx+1} 변환 중 오류: {e}")
        import traceback
        traceback.print_exc()

# PPTX 파일 저장
output_pptx = "merged_slides.pptx"
prs.save(output_pptx)

print(f"\n✓ PPTX 변환 완료: {output_pptx}")
print(f"  총 {len(slides)}개 슬라이드가 포함되었습니다.")
print(f"  모든 텍스트와 객체가 개별적으로 편집 가능합니다!")

